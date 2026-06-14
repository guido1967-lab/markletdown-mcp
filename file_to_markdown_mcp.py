#!/usr/bin/env python3
"""
File-to-Markdown MCP Server for Claude Desktop
Converts PDFs, Word, Excel, PowerPoint, images, etc. to Markdown
"""

import os
import sys
import json
import subprocess
import tempfile
from pathlib import Path
from typing import Any

# Try importing optional dependencies
try:
    from pdf2image import convert_from_path
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    import pytesseract
    from PIL import Image
    IMAGE_SUPPORT = True
except ImportError:
    IMAGE_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    import openpyxl
    from openpyxl import load_workbook
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False


class FileToMarkdownConverter:
    """Convert various file formats to Markdown"""
    
    SUPPORTED_FORMATS = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document',
        '.doc': 'Word Document (97-2003)',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel Spreadsheet (97-2003)',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint Presentation (97-2003)',
        '.jpg': 'JPEG Image',
        '.jpeg': 'JPEG Image',
        '.png': 'PNG Image',
        '.gif': 'GIF Image',
        '.bmp': 'Bitmap Image',
        '.tiff': 'TIFF Image',
        '.txt': 'Plain Text',
        '.csv': 'CSV File',
        '.html': 'HTML File',
        '.json': 'JSON File',
    }
    
    def __init__(self):
        self.temp_dir = tempfile.gettempdir()
    
    def convert_file(self, file_path: str) -> dict:
        """
        Convert a file to Markdown format
        Returns: {'status': 'success'|'error', 'markdown': str, 'message': str}
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            return {
                'status': 'error',
                'markdown': '',
                'message': f'File not found: {file_path}'
            }
        
        file_ext = file_path.suffix.lower()
        
        try:
            if file_ext == '.pdf':
                return self._convert_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._convert_word(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._convert_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._convert_pptx(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return self._convert_image(file_path)
            elif file_ext == '.txt':
                return self._convert_text(file_path)
            elif file_ext == '.csv':
                return self._convert_csv(file_path)
            elif file_ext == '.html':
                return self._convert_html(file_path)
            elif file_ext == '.json':
                return self._convert_json(file_path)
            else:
                return {
                    'status': 'error',
                    'markdown': '',
                    'message': f'Unsupported file format: {file_ext}\nSupported: {", ".join(self.SUPPORTED_FORMATS.keys())}'
                }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting file: {str(e)}'
            }
    
    def _convert_pdf(self, file_path: Path) -> dict:
        """Convert PDF to Markdown"""
        try:
            # Use pandoc for PDF conversion
            result = subprocess.run(
                ['pandoc', str(file_path), '-t', 'markdown'],
                capture_output=True,
                text=True,
                check=True
            )
            return {
                'status': 'success',
                'markdown': result.stdout,
                'message': f'Successfully converted PDF: {file_path.name}'
            }
        except subprocess.CalledProcessError as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Pandoc error: {e.stderr}'
            }
    
    def _convert_word(self, file_path: Path) -> dict:
        """Convert Word document to Markdown"""
        try:
            result = subprocess.run(
                ['pandoc', str(file_path), '-t', 'markdown'],
                capture_output=True,
                text=True,
                check=True
            )
            return {
                'status': 'success',
                'markdown': result.stdout,
                'message': f'Successfully converted Word document: {file_path.name}'
            }
        except subprocess.CalledProcessError as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Pandoc error: {e.stderr}'
            }
    
    def _convert_excel(self, file_path: Path) -> dict:
        """Convert Excel spreadsheet to Markdown tables"""
        try:
            if not EXCEL_SUPPORT:
                return {
                    'status': 'error',
                    'markdown': '',
                    'message': 'openpyxl not installed. Run: pip3 install openpyxl'
                }
            
            workbook = load_workbook(file_path)
            markdown = f"# {file_path.stem}\n\n"
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                markdown += f"## Sheet: {sheet_name}\n\n"
                
                # Get dimensions
                rows = list(sheet.iter_rows(values_only=True))
                
                if rows:
                    # Create header
                    headers = rows[0]
                    markdown += "| " + " | ".join(str(h or "") for h in headers) + " |\n"
                    markdown += "|" + "|".join(["---"] * len(headers)) + "|\n"
                    
                    # Add data rows
                    for row in rows[1:]:
                        markdown += "| " + " | ".join(str(cell or "") for cell in row) + " |\n"
                    
                    markdown += "\n"
            
            return {
                'status': 'success',
                'markdown': markdown,
                'message': f'Successfully converted Excel file: {file_path.name}'
            }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting Excel: {str(e)}'
            }
    
    def _convert_pptx(self, file_path: Path) -> dict:
        """Convert PowerPoint to Markdown"""
        try:
            if not PPTX_SUPPORT:
                return {
                    'status': 'error',
                    'markdown': '',
                    'message': 'python-pptx not installed. Run: pip3 install python-pptx'
                }
            
            prs = Presentation(file_path)
            markdown = f"# {file_path.stem}\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                markdown += f"## Slide {i}\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        markdown += f"{shape.text}\n\n"
                    elif hasattr(shape, "table"):
                        # Handle tables in slides
                        table = shape.table
                        rows = table.rows
                        cols = table.columns
                        
                        for row_idx, row in enumerate(rows):
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text)
                            markdown += "| " + " | ".join(row_data) + " |\n"
                            
                            if row_idx == 0:
                                markdown += "|" + "|".join(["---"] * len(cols)) + "|\n"
                        markdown += "\n"
            
            return {
                'status': 'success',
                'markdown': markdown,
                'message': f'Successfully converted PowerPoint: {file_path.name}'
            }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting PowerPoint: {str(e)}'
            }
    
    def _convert_image(self, file_path: Path) -> dict:
        """Convert image to Markdown with OCR"""
        try:
            if not IMAGE_SUPPORT:
                return {
                    'status': 'error',
                    'markdown': '',
                    'message': 'Pillow or pytesseract not installed. Run: pip3 install pillow pytesseract'
                }
            
            # Try OCR first
            try:
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                if text.strip():
                    markdown = f"# Image: {file_path.stem}\n\n## Extracted Text\n\n{text}\n\n"
                    markdown += f"![{file_path.name}]({file_path.name})\n"
                    return {
                        'status': 'success',
                        'markdown': markdown,
                        'message': f'Successfully extracted text from image: {file_path.name}'
                    }
            except:
                pass
            
            # Fallback: just reference the image
            markdown = f"# Image: {file_path.stem}\n\n![{file_path.name}]({file_path.name})\n"
            return {
                'status': 'success',
                'markdown': markdown,
                'message': f'Image added to Markdown: {file_path.name} (OCR not available)'
            }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting image: {str(e)}'
            }
    
    def _convert_text(self, file_path: Path) -> dict:
        """Convert plain text file to Markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            markdown = f"# {file_path.stem}\n\n{content}\n"
            return {
                'status': 'success',
                'markdown': markdown,
                'message': f'Successfully converted text file: {file_path.name}'
            }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting text file: {str(e)}'
            }
    
    def _convert_csv(self, file_path: Path) -> dict:
        """Convert CSV to Markdown table"""
        try:
            import csv
            
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)
            
            if not rows:
                return {
                    'status': 'error',
                    'markdown': '',
                    'message': 'CSV file is empty'
                }
            
            markdown = f"# {file_path.stem}\n\n"
            
            # Create header
            headers = rows[0]
            markdown += "| " + " | ".join(headers) + " |\n"
            markdown += "|" + "|".join(["---"] * len(headers)) + "|\n"
            
            # Add data rows
            for row in rows[1:]:
                # Pad row if needed
                row = row + [''] * (len(headers) - len(row))
                markdown += "| " + " | ".join(row[:len(headers)]) + " |\n"
            
            return {
                'status': 'success',
                'markdown': markdown,
                'message': f'Successfully converted CSV file: {file_path.name}'
            }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting CSV: {str(e)}'
            }
    
    def _convert_html(self, file_path: Path) -> dict:
        """Convert HTML to Markdown"""
        try:
            result = subprocess.run(
                ['pandoc', str(file_path), '-f', 'html', '-t', 'markdown'],
                capture_output=True,
                text=True,
                check=True
            )
            return {
                'status': 'success',
                'markdown': result.stdout,
                'message': f'Successfully converted HTML file: {file_path.name}'
            }
        except subprocess.CalledProcessError as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Pandoc error: {e.stderr}'
            }
    
    def _convert_json(self, file_path: Path) -> dict:
        """Convert JSON to Markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            markdown = f"# {file_path.stem}\n\n```json\n{json.dumps(data, indent=2)}\n```\n"
            return {
                'status': 'success',
                'markdown': markdown,
                'message': f'Successfully converted JSON file: {file_path.name}'
            }
        except Exception as e:
            return {
                'status': 'error',
                'markdown': '',
                'message': f'Error converting JSON: {str(e)}'
            }


def main():
    """Main entry point"""
    converter = FileToMarkdownConverter()
    
    if len(sys.argv) < 2:
        print("Usage: python3 file_to_markdown_mcp.py <file_path>")
        print("\nSupported formats:")
        for fmt, desc in converter.SUPPORTED_FORMATS.items():
            print(f"  {fmt}: {desc}")
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = converter.convert_file(file_path)
    
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
